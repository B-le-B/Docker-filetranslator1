import os
import logging
import shutil
import tempfile
import re
import time
import hashlib
import threading
import copy
from typing import Optional, List, Dict, Any, Tuple, Union
from enum import Enum
from dataclasses import dataclass, field
from collections import OrderedDict
from contextlib import contextmanager
import datetime
import decimal
from pathlib import Path
from tqdm import tqdm
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame
from pptx.table import Table as PPTTable
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches
from pptx.shapes.placeholder import (
    LayoutPlaceholder, SlidePlaceholder, MasterPlaceholder, 
    NotesSlidePlaceholder, BasePlaceholder, PlaceholderPicture
)
from concurrent.futures import ThreadPoolExecutor, as_completed


def is_placeholder(shape):
    """检查shape是否是placeholder"""
    return isinstance(shape, (LayoutPlaceholder, SlidePlaceholder, 
                            MasterPlaceholder, NotesSlidePlaceholder, BasePlaceholder))

class NoDebugFilter(logging.Filter):
    def filter(self, record):
        if record.levelno >= logging.WARNING:
            return True
        if record.name == __name__ or record.name.startswith('__main__'):
            return record.levelno >= logging.INFO
        return False

root_logger = logging.getLogger()
root_logger.addFilter(NoDebugFilter())

logger = logging.getLogger(__name__)

# 预编译正则表达式
NUMBERED_LINE_PATTERN = re.compile(r'^\[(\d+)\]\s*(.*)')
LIST_MARKER_PATTERN = re.compile(r'^([\d\w]+[\.\)]\s*|[•·▪▫◦‣⁃]\s*)')
WHITESPACE_PATTERN = re.compile(r'\s+')

# 重用现有的枚举和数据类
class FailureReason(Enum):
    """失败原因枚举"""
    TIMEOUT = "timeout"
    API_ERROR = "api_error"
    PARSE_ERROR = "parse_error"
    NOT_TRANSLATED = "not_translated"
    EMPTY_RESPONSE = "empty_response"
    CONNECTION_ERROR = "connection_error"
    BATCH_FAILURE = "batch_failure"

@dataclass
class RunFormatInfo:
    """Run级别格式信息"""
    text: str = ""
    font_name: Optional[str] = None
    font_size: Optional[int] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    font_color_rgb: Optional[Tuple[int, int, int]] = None
    highlight_color_rgb: Optional[Tuple[int, int, int]] = None
    start_pos: int = 0
    end_pos: int = 0

@dataclass
class ParagraphFormatInfo:
    """段落级别格式信息"""
    alignment: Optional[int] = None
    line_spacing: Optional[float] = None
    space_before: Optional[int] = None
    space_after: Optional[int] = None
    left_indent: Optional[int] = None
    first_line_indent: Optional[int] = None
    bullet_level: Optional[int] = None
    bullet_font: Optional[str] = None
    bullet_char: Optional[str] = None
    runs_format: List[RunFormatInfo] = field(default_factory=list)

@dataclass
class FailedTask:
    """统一的失败任务数据类"""
    original_text: str
    original_index: int
    element: Any  # PPTElement对象
    failure_reason: FailureReason
    retry_count: int = 0
    error_message: str = ""
    is_serious: bool = True
    def __post_init__(self):
        self.failure_timestamp = time.time()

@dataclass
class PPTElement:
    """PPT可翻译元素，保持幻灯片完整性"""
    full_text: str
    slide_index: int
    shape: BaseShape
    shape_type: str  # 'textbox', 'table_cell', 'placeholder', 'chart', 'smartart'
    shape_info: Dict[str, Any]  # 形状的位置、大小、格式信息
    text_container: Any  # TextFrame 或 Cell 对象
    paragraph_index: Optional[int] = None  # 如果是段落级别的翻译
    cell_info: Optional[Dict[str, Any]] = None  # 表格单元格信息
    location: str = ""  # 位置描述
    format_info: ParagraphFormatInfo = None  # 增强的格式信息
    element_type: str = 'ppt_element'

    def __post_init__(self):
        if self.format_info is None:
            self.format_info = ParagraphFormatInfo()

# 格式处理增强函数
def _safe_get_color_rgb(color_obj) -> Optional[Tuple[int, int, int]]:
    """安全获取颜色RGB值"""
    try:
        if color_obj and hasattr(color_obj, 'rgb') and color_obj.rgb:
            rgb = color_obj.rgb
            return (rgb.r, rgb.g, rgb.b)
        return None
    except Exception:
        return None

def _safe_set_color_rgb(color_obj, rgb_tuple: Tuple[int, int, int]):
    """安全设置颜色RGB值"""
    try:
        if color_obj and rgb_tuple and len(rgb_tuple) == 3:
            color_obj.rgb = RGBColor(*rgb_tuple)
    except Exception as e:
        logger.debug(f"Failed to set color RGB: {e}")

def _extract_run_format_info(run, start_pos: int = 0) -> RunFormatInfo:
    """提取单个run的完整格式信息"""
    try:
        run_info = RunFormatInfo(
            text=run.text,
            start_pos=start_pos,
            end_pos=start_pos + len(run.text)
        )
        
        # 字体信息
        if hasattr(run, 'font'):
            font = run.font
            try:
                run_info.font_name = font.name
            except:
                pass
            try:
                run_info.font_size = int(font.size.pt) if font.size else None
            except:
                pass
            try:
                run_info.bold = font.bold
            except:
                pass
            try:
                run_info.italic = font.italic
            except:
                pass
            try:
                run_info.underline = font.underline
            except:
                pass
            
            # 字体颜色
            try:
                run_info.font_color_rgb = _safe_get_color_rgb(font.color)
            except:
                pass
            
            # 高亮颜色
            try:
                if hasattr(font, 'highlight_color'):
                    run_info.highlight_color_rgb = _safe_get_color_rgb(font.highlight_color)
            except:
                pass
        
        return run_info
    except Exception as e:
        logger.debug(f"Error extracting run format: {e}")
        return RunFormatInfo(text=run.text if hasattr(run, 'text') else "", start_pos=start_pos)

def _extract_paragraph_format_info_enhanced(paragraph) -> ParagraphFormatInfo:
    """提取段落的完整格式信息（增强版）"""
    try:
        para_info = ParagraphFormatInfo()
        
        # 段落级格式
        try:
            para_info.alignment = paragraph.alignment
        except:
            pass
        
        try:
            para_info.line_spacing = float(paragraph.line_spacing) if paragraph.line_spacing else None
        except:
            pass
        
        try:
            para_info.space_before = int(paragraph.space_before.pt) if paragraph.space_before else None
        except:
            pass
        
        try:
            para_info.space_after = int(paragraph.space_after.pt) if paragraph.space_after else None
        except:
            pass
        
        try:
            para_info.left_indent = int(paragraph.left_indent.pt) if paragraph.left_indent else None
        except:
            pass
        
        try:
            para_info.first_line_indent = int(paragraph.first_line_indent.pt) if paragraph.first_line_indent else None
        except:
            pass
        
        # 项目符号信息
        try:
            para_info.bullet_level = paragraph.level
        except:
            pass
        
        try:
            if hasattr(paragraph, 'font'):
                para_info.bullet_font = paragraph.font.name
        except:
            pass
        
        # 提取所有run的格式信息
        current_pos = 0
        for run in paragraph.runs:
            run_info = _extract_run_format_info(run, current_pos)
            para_info.runs_format.append(run_info)
            current_pos = run_info.end_pos
        
        return para_info
    except Exception as e:
        logger.debug(f"Error extracting paragraph format: {e}")
        return ParagraphFormatInfo()

def _apply_run_format_info(run, run_info: RunFormatInfo):
    """应用run格式信息"""
    try:
        if not hasattr(run, 'font'):
            return
        
        font = run.font
        
        # 基础字体属性
        if run_info.font_name:
            try:
                font.name = run_info.font_name
            except:
                pass
        
        if run_info.font_size:
            try:
                font.size = Pt(run_info.font_size)
            except:
                pass
        
        if run_info.bold is not None:
            try:
                font.bold = run_info.bold
            except:
                pass
        
        if run_info.italic is not None:
            try:
                font.italic = run_info.italic
            except:
                pass
        
        if run_info.underline is not None:
            try:
                font.underline = run_info.underline
            except:
                pass
        
        # 颜色信息
        if run_info.font_color_rgb:
            try:
                _safe_set_color_rgb(font.color, run_info.font_color_rgb)
            except:
                pass
        
        if run_info.highlight_color_rgb:
            try:
                if hasattr(font, 'highlight_color'):
                    _safe_set_color_rgb(font.highlight_color, run_info.highlight_color_rgb)
            except:
                pass
        
    except Exception as e:
        logger.debug(f"Error applying run format: {e}")

def _apply_paragraph_format_info(paragraph, para_info: ParagraphFormatInfo):
    """应用段落格式信息"""
    try:
        # 段落级格式
        if para_info.alignment is not None:
            try:
                paragraph.alignment = para_info.alignment
            except:
                pass
        
        if para_info.line_spacing is not None:
            try:
                paragraph.line_spacing = para_info.line_spacing
            except:
                pass
        
        if para_info.space_before is not None:
            try:
                paragraph.space_before = Pt(para_info.space_before)
            except:
                pass
        
        if para_info.space_after is not None:
            try:
                paragraph.space_after = Pt(para_info.space_after)
            except:
                pass
        
        if para_info.left_indent is not None:
            try:
                paragraph.left_indent = Pt(para_info.left_indent)
            except:
                pass
        
        if para_info.first_line_indent is not None:
            try:
                paragraph.first_line_indent = Pt(para_info.first_line_indent)
            except:
                pass
        
        # 项目符号级别
        if para_info.bullet_level is not None:
            try:
                paragraph.level = para_info.bullet_level
            except:
                pass
        
    except Exception as e:
        logger.debug(f"Error applying paragraph format: {e}")

def _rebuild_paragraph_with_format(paragraph, translated_text: str, para_info: ParagraphFormatInfo):
    """重构段落并保持格式（核心格式保持函数）"""
    try:
        # 保存段落级格式设置
        saved_alignment = para_info.alignment
        saved_level = para_info.bullet_level
        
        # 清空段落内容但保持段落对象
        paragraph.clear()
        
        # 恢复段落级格式
        _apply_paragraph_format_info(paragraph, para_info)
        
        # 智能分配格式到新文本
        if para_info.runs_format:
            # 根据原有run格式分配新文本
            if len(para_info.runs_format) == 1:
                # 单一格式，直接应用
                run = paragraph.add_run()
                run.text = translated_text
                _apply_run_format_info(run, para_info.runs_format[0])
            else:
                # 多种格式，智能分配
                _distribute_text_with_mixed_format(paragraph, translated_text, para_info.runs_format)
        else:
            # 没有格式信息，使用默认
            run = paragraph.add_run()
            run.text = translated_text
        
    except Exception as e:
        logger.debug(f"Error rebuilding paragraph with format: {e}")
        # 失败时的简单恢复
        paragraph.clear()
        run = paragraph.add_run()
        run.text = translated_text

def _distribute_text_with_mixed_format(paragraph, translated_text: str, runs_format: List[RunFormatInfo]):
    """为混合格式文本智能分配格式"""
    try:
        # 策略1：按长度比例分配
        total_original_length = sum(len(run_info.text) for run_info in runs_format)
        if total_original_length == 0:
            # 没有原始文本，使用第一个格式
            run = paragraph.add_run()
            run.text = translated_text
            _apply_run_format_info(run, runs_format[0])
            return
        
        current_pos = 0
        for run_info in runs_format:
            if current_pos >= len(translated_text):
                break
            
            # 计算当前run应该分配的长度
            original_ratio = len(run_info.text) / total_original_length
            allocated_length = max(1, int(len(translated_text) * original_ratio))
            
            # 确保不超出边界
            end_pos = min(current_pos + allocated_length, len(translated_text))
            
            # 创建run并应用格式
            if end_pos > current_pos:
                run = paragraph.add_run()
                run.text = translated_text[current_pos:end_pos]
                _apply_run_format_info(run, run_info)
                current_pos = end_pos
        
        # 处理剩余文本
        if current_pos < len(translated_text):
            run = paragraph.add_run()
            run.text = translated_text[current_pos:]
            # 使用最后一个格式
            if runs_format:
                _apply_run_format_info(run, runs_format[-1])
        
    except Exception as e:
        logger.debug(f"Error distributing text with mixed format: {e}")
        # 失败时使用简单策略
        run = paragraph.add_run()
        run.text = translated_text
        if runs_format:
            _apply_run_format_info(run, runs_format[0])

# 重用现有的工具函数
def _prepare_prompt_config(prompt_config: Optional[Dict[str, Any]], kwargs: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """准备和标准化prompt配置，与前端格式兼容"""
    try:
        has_prompt_config = prompt_config and isinstance(prompt_config, dict) and prompt_config
        has_kwargs_config = any(k in kwargs for k in ['preserve_terms', 'glossary', 'additional_context', 'prompt_template', 'custom_prompt'])
        if not has_prompt_config and not has_kwargs_config:
            logger.debug("No valid prompt configuration found, returning None")
            return None
        config = {}
        if has_prompt_config:
            config.update(prompt_config)
        for key in ['preserve_terms', 'glossary', 'additional_context', 'prompt_template', 'custom_prompt']:
            if key in kwargs and kwargs[key] is not None:
                config[key] = kwargs[key]
        if not config:
            logger.debug("Config is empty after processing, returning None")
            return None
        normalized_config = _normalize_prompt_config(config)
        if not normalized_config or normalized_config.get('mode') == 'none':
            logger.debug("Normalized config is invalid, returning None")
            return None
        logger.debug(f"Prepared prompt config for advanced PPTX translation: {normalized_config}")
        return normalized_config
    except Exception as e:
        logger.warning(f"Error preparing prompt config: {e}")
        return None

def _normalize_prompt_config(config: Dict[str, Any]) -> Dict[str, Any]:
    """标准化prompt配置格式"""
    try:
        if not config or not isinstance(config, dict):
            return {'mode': 'none'}
        normalized = config.copy()
        if 'mode' not in normalized:
            if 'custom_prompt' in normalized and normalized['custom_prompt']:
                normalized['mode'] = 'custom'
            elif 'prompt_template' in normalized or 'professional_domain' in normalized:
                normalized['mode'] = 'professional'
            elif any(k in normalized for k in ['preserve_terms', 'glossary', 'additional_context']):
                normalized['mode'] = 'general'
            else:
                normalized['mode'] = 'none'
        # 处理保留术语
        preserve_terms = normalized.get('preserve_terms')
        if preserve_terms:
            try:
                if isinstance(preserve_terms, str):
                    terms_list = [term.strip() for term in preserve_terms.split(',') if term.strip()]
                    normalized['preserve_terms'] = terms_list
                elif isinstance(preserve_terms, list):
                    normalized['preserve_terms'] = [str(term).strip() for term in preserve_terms if str(term).strip()]
                else:
                    normalized.pop('preserve_terms', None)
            except Exception as e:
                logger.warning(f"Error processing preserve_terms: {e}")
                normalized.pop('preserve_terms', None)
        # 处理术语表
        glossary = normalized.get('glossary')
        if glossary:
            if not isinstance(glossary, dict):
                logger.warning(f"Glossary should be a dictionary, got {type(glossary)}, ignoring")
                normalized.pop('glossary', None)
        # 处理自定义prompt
        if normalized.get('mode') == 'custom':
            custom_prompt = normalized.get('custom_prompt', {})
            if not custom_prompt or not isinstance(custom_prompt, dict):
                system_prompt = normalized.get('custom_system_prompt', normalized.get('system'))
                user_prompt = normalized.get('custom_user_prompt', normalized.get('user'))
                if system_prompt:
                    normalized['custom_prompt'] = {
                        'system': str(system_prompt),
                        'user': str(user_prompt) if user_prompt else 'Please translate the following content to {target_lang}:\n\n{content}'
                    }
                else:
                    logger.warning("Custom mode selected but no valid custom prompt provided, falling back to general mode")
                    normalized['mode'] = 'general'
        # 处理专业模板
        if normalized.get('mode') == 'professional':
            domain = normalized.get('professional_domain', normalized.get('prompt_template', 'academic'))
            normalized['prompt_template'] = str(domain)
        return normalized
    except Exception as e:
        logger.warning(f"Error normalizing prompt config: {e}")
        return {'mode': 'none'}

def _get_batch_settings_from_config(prompt_config: Optional[Dict[str, Any]], kwargs: Dict[str, Any]) -> Dict[str, int]:
    """从配置中获取批处理设置"""
    settings = {
        'batch_size': 50,  # PPT通常文本更少，减小批次
        'max_chars': 6000,  # PPT文本通常更短
        'max_workers': 4,   # PPT处理相对简单，减少worker
        'retry_max_workers': 3,
        'translation_timeout': 45,
        'max_retries': 6,
        'large_text_threshold': 30,  # PPT大段文本阈值更低
        'retry_failure_threshold': 0.0,
        'non_ascii_threshold': 0.0
    }
    try:
        if 'batch_size' in kwargs and isinstance(kwargs['batch_size'], (int, float)):
            settings['batch_size'] = max(3, min(int(kwargs['batch_size']), 100))
        if 'max_chunk_size' in kwargs and isinstance(kwargs['max_chunk_size'], (int, float)):
            settings['max_chars'] = max(500, min(int(kwargs['max_chunk_size']), 20000))
        if 'max_workers' in kwargs and isinstance(kwargs['max_workers'], (int, float)):
            settings['max_workers'] = max(1, min(int(kwargs['max_workers']), 8))
        if 'retry_max_workers' in kwargs and isinstance(kwargs['retry_max_workers'], (int, float)):
            settings['retry_max_workers'] = max(1, min(int(kwargs['retry_max_workers']), 6))
        if 'translation_timeout' in kwargs and isinstance(kwargs['translation_timeout'], (int, float)):
            settings['translation_timeout'] = max(20, min(int(kwargs['translation_timeout']), 180))
        if 'max_retries' in kwargs and isinstance(kwargs['max_retries'], (int, float)):
            settings['max_retries'] = max(2, min(int(kwargs['max_retries']), 10))
        if prompt_config and isinstance(prompt_config, dict):
            max_units = prompt_config.get('max_units_per_chunk')
            if max_units and isinstance(max_units, (int, float)):
                settings['batch_size'] = max(3, min(int(max_units), 50))
            max_chars = prompt_config.get('max_chars_per_chunk')
            if max_chars and isinstance(max_chars, (int, float)):
                settings['max_chars'] = max(500, min(int(max_chars), 20000))
    except Exception as e:
        logger.warning(f"Error processing batch settings: {e}")
    return settings

# 重用现有的缓存系统
class SmartCache:
    """智能LRU缓存，支持prompt配置差异化和失败缓存清理"""
    def __init__(self, max_size: int = 800):  # PPT缓存稍小
        self._cache = OrderedDict()
        self.max_size = max_size
        self._lock = threading.RLock()
        self._hits = 0
        self._misses = 0
        self._clears = 0
    def get(self, key: str) -> Optional[str]:
        if not key or not isinstance(key, str):
            return None
        with self._lock:
            if key in self._cache:
                value = self._cache.pop(key)
                self._cache[key] = value
                self._hits += 1
                return value
            self._misses += 1
            return None
    def put(self, key: str, value: str):
        if not key or not isinstance(key, str) or not isinstance(value, str):
            return
        with self._lock:
            if key in self._cache:
                self._cache.pop(key)
            elif len(self._cache) >= self.max_size:
                self._cache.popitem(last=False)
            self._cache[key] = value
    def remove(self, key: str) -> bool:
        if not key or not isinstance(key, str):
            return False
        with self._lock:
            if key in self._cache:
                self._cache.pop(key)
                self._clears += 1
                return True
            return False
    def clear_pattern(self, pattern: str) -> int:
        if not pattern:
            return 0
        with self._lock:
            keys_to_remove = [key for key in self._cache.keys() if pattern in key]
            for key in keys_to_remove:
                self._cache.pop(key)
            self._clears += len(keys_to_remove)
            return len(keys_to_remove)
    def clear(self):
        with self._lock:
            cleared_count = len(self._cache)
            self._cache.clear()
            self._hits = 0
            self._misses = 0
            self._clears += cleared_count
    @property
    def stats(self) -> Dict[str, Any]:
        with self._lock:
            total = self._hits + self._misses
            return {
                'size': len(self._cache),
                'hits': self._hits,
                'misses': self._misses,
                'clears': self._clears,
                'hit_rate': self._hits / total if total > 0 else 0
            }

# 重用翻译验证器
class TranslationValidator:
    """翻译完整性验证器"""
    ERROR_KEYWORDS = [
        'timeout', 'readtimeout', 'connecttimeout', 'httptimeout',
        'network error', 'connection error', 'api error', 'service error',
        'translation failed', 'service unavailable', 'request failed',
        'server error', 'bad gateway', 'gateway timeout',
        '超时', '网络错误', '连接错误', '服务错误', '翻译失败',
        '服务不可用', '请求失败', '服务器错误'
    ]
    @staticmethod
    def is_error_message(text: str) -> bool:
        if not text:
            return True
        text_lower = text.lower()
        return any(keyword.lower() in text_lower for keyword in TranslationValidator.ERROR_KEYWORDS)
    @staticmethod
    def is_serious_failure(original_text: str, translated_text: str, large_text_threshold: int = 30, from_cache: bool = False) -> Tuple[bool, str]:
        """PPT版本的失败判断，阈值调整为30"""
        if from_cache:
            return False, "缓存结果"
        if not translated_text or translated_text.strip() == "":
            if len(original_text) > large_text_threshold:
                return True, f"大段文字未翻译（{len(original_text)}字符）"
            else:
                return False, f"短文本未翻译（{len(original_text)}字符，可能正常）"
        if TranslationValidator.is_error_message(translated_text):
            return True, f"翻译服务错误: {translated_text[:50]}..."
        if original_text.strip() == translated_text.strip():
            if len(original_text) > large_text_threshold:
                return True, f"长文本未翻译（{len(original_text)}字符）"
            else:
                return False, f"短文本与原文相同（{len(original_text)}字符，正常情况）"
        if from_cache and len(original_text) > large_text_threshold:
            if original_text.strip() == translated_text.strip():
                return True, f"缓存中的长文本未翻译（{len(original_text)}字符）"
        return False, "翻译成功"

class AdvancedPptxTranslator:
    """高级PPTX翻译器 - 格式保持增强版"""
    def __init__(
        self, 
        translator, 
        batch_size: int = 50,
        max_chars: int = 6000,
        max_workers: int = 4,
        retry_max_workers: int = 3,
        prompt_config: Optional[Dict[str, Any]] = None,
        reference_doc: Optional[str] = None,
        translation_timeout: int = 45,
        max_retries: int = 6,
        large_text_threshold: int = 30,
        retry_failure_threshold: float = 0.0,
        non_ascii_threshold: float = 0.0,
        translate_notes: bool = False,  # PPT特有：是否翻译备注
        translate_slide_titles: bool = True,  # PPT特有：是否翻译幻灯片标题
        preserve_formatting: bool = True,  # 新增：是否保持格式
        **kwargs
    ):
        self.translator = translator
        self.batch_size = batch_size
        self.max_chars = max_chars
        self.max_workers = max_workers
        self.retry_max_workers = retry_max_workers
        self.translation_timeout = translation_timeout
        self.max_retries = max_retries
        self.large_text_threshold = large_text_threshold
        self.retry_failure_threshold = retry_failure_threshold
        self.non_ascii_threshold = non_ascii_threshold
        # PPT特有设置
        self.translate_notes = translate_notes
        self.translate_slide_titles = translate_slide_titles
        self.preserve_formatting = preserve_formatting  # 新增格式保持开关
        
        # 验证worker配置
        if self.retry_max_workers < 1:
            self.retry_max_workers = 1
            logger.warning("retry_max_workers不能小于1，已重置为1")
        if self.retry_max_workers > 6:
            logger.warning(f"retry_max_workers={self.retry_max_workers}可能过高，建议不超过6")
        
        # Worker资源管理
        total_max_workers = 10
        if self.max_workers + self.retry_max_workers > total_max_workers:
            self.max_workers = max(1, total_max_workers - self.retry_max_workers)
            logger.info(f"调整主翻译workers为{self.max_workers}，为重试预留{self.retry_max_workers}个workers")
        
        self.cache = SmartCache(800)
        self.source_lang = None
        self.reference_doc = reference_doc
        self._config_lock = threading.RLock()
        self._cached_config_hash = None
        
        # 失败任务追踪
        self.failed_tasks: List[FailedTask] = []
        self.failed_tasks_lock = threading.Lock()
        
        # PPT调整的重试策略 - 更保守的批次大小
        self.retry_batch_sizes = [10, 5, 2, 1, 1, 1]
        self.retry_delays = [1, 2, 4, 6, 8, 10]
        
        # 处理prompt配置
        try:
            self.effective_prompt_config = _prepare_prompt_config(prompt_config, kwargs)
            self.original_translator_config = None
            if self.effective_prompt_config is None:
                logger.debug("No effective prompt config, using default mode")
                self.effective_prompt_config = {'mode': 'none'}
        except Exception as e:
            logger.warning(f"Error initializing prompt config: {e}")
            self.effective_prompt_config = {'mode': 'none'}
        
        # 从配置获取批处理设置
        try:
            batch_settings = _get_batch_settings_from_config(self.effective_prompt_config, kwargs)
            self.batch_size = batch_settings['batch_size']
            self.max_chars = batch_settings['max_chars']
            self.max_workers = min(self.max_workers, batch_settings['max_workers'])
            self.retry_max_workers = min(self.retry_max_workers, batch_settings['retry_max_workers'])
            self.translation_timeout = batch_settings['translation_timeout']
            self.max_retries = batch_settings['max_retries']
            self.large_text_threshold = batch_settings['large_text_threshold']
            self.retry_failure_threshold = batch_settings['retry_failure_threshold']
            self.non_ascii_threshold = batch_settings['non_ascii_threshold']
        except Exception as e:
            logger.warning(f"Error getting batch settings: {e}")
        
        if self.effective_prompt_config and self.effective_prompt_config.get('mode') != 'none':
            logger.info(f"AdvancedPptxTranslator initialized with prompt config: mode={self.effective_prompt_config.get('mode')}")
        
        logger.info(f"格式保持模式: {'启用' if self.preserve_formatting else '禁用'}")
        
        # PPT特有统计
        self.stats = {
            'total_slides': 0,
            'total_elements': 0,
            'translated_elements': 0,
            'skipped_elements': 0,
            'total_chars': 0,
            'total_batches': 0,
            'tables': 0,
            'text_boxes': 0,
            'placeholders': 0,
            'notes_pages': 0,
            'charts': 0,
            'smartart': 0,
            'processing_time': 0,
            'prompt_mode': self.effective_prompt_config.get('mode', 'none'),
            'api_calls': 0,
            'cache_savings': 0,
            'template_applied': False,
            'serious_failures': 0,
            'minor_issues': 0,
            'retry_attempts': 0,
            'final_failures': 0,
            'cache_clears': 0,
            'final_rescues': 0,
            'retry_workers_used': 0,
            'concurrent_retry_batches': 0,
            'format_preserved': 0,  # 新增：格式保持统计
            'format_fallback': 0,   # 新增：格式恢复失败统计
        }

    @contextmanager
    def _translator_config_context(self):
        """线程安全的配置上下文管理器"""
        try:
            self._apply_prompt_config_to_translator()
            yield
        except Exception as e:
            logger.error(f"Error in translator config context: {e}")
            raise
        finally:
            self._restore_translator_config()

    def _apply_prompt_config_to_translator(self):
        """应用prompt配置到翻译器"""
        if (self.effective_prompt_config and 
            self.effective_prompt_config.get('mode') != 'none' and 
            hasattr(self.translator, 'set_prompt_config')):
            try:
                with self._config_lock:
                    self.original_translator_config = getattr(self.translator, 'prompt_config', None)
                    config_copy = copy.deepcopy(self.effective_prompt_config)
                    self.translator.set_prompt_config(config_copy)
                    logger.info("Applied prompt config to translator in AdvancedPptxTranslator")
                    return True
            except Exception as e:
                logger.warning(f"Failed to apply prompt config to translator: {e}")
        return False

    def _restore_translator_config(self):
        """恢复翻译器的原始配置"""
        if (self.effective_prompt_config and 
            self.effective_prompt_config.get('mode') != 'none' and 
            hasattr(self.translator, 'set_prompt_config')):
            try:
                with self._config_lock:
                    if self.original_translator_config is not None:
                        self.translator.set_prompt_config(self.original_translator_config)
                    else:
                        if hasattr(self.translator, 'prompt_config'):
                            self.translator.prompt_config = None
                    logger.debug("Restored translator config in AdvancedPptxTranslator")
            except Exception as e:
                logger.warning(f"Failed to restore translator config: {e}")

    def _get_config_safe(self) -> Optional[Dict[str, Any]]:
        """线程安全地获取配置副本"""
        with self._config_lock:
            return copy.deepcopy(self.effective_prompt_config) if self.effective_prompt_config else None


# pptx_full_translator.py (第二批 - 格式保持增强版)

    def _clear_failed_task_cache(self, task: FailedTask, target_lang: str, source_lang: Optional[str]):
        """清理失败任务的缓存"""
        try:
            local_prompt_config = self._get_config_safe()
            cache_key = self._get_cache_key(task.original_text, target_lang, source_lang, local_prompt_config)
            if self.cache.remove(cache_key):
                self.stats['cache_clears'] += 1
                logger.debug(f"清理失败任务缓存: {task.original_text[:50]}...")
        except Exception as e:
            logger.warning(f"Error clearing failed task cache: {e}")

    def _should_trigger_retry(self, total_processed: int) -> bool:
        """判断是否应该触发重试"""
        with self.failed_tasks_lock:
            serious_failures = [task for task in self.failed_tasks if task.is_serious]
            if not serious_failures:
                return False
            failure_rate = len(serious_failures) / max(1, total_processed)
            should_retry = failure_rate >= self.retry_failure_threshold
            logger.info(f"严重失败: {len(serious_failures)}/{total_processed} ({failure_rate:.1%}), "
                       f"阈值: {self.retry_failure_threshold:.1%}, 是否重试: {should_retry}")
            return should_retry

    def _add_failed_task(self, element: PPTElement, original_index: int, 
                        reason: str, from_cache: bool = False):
        """添加失败任务，智能判断是否为严重失败"""
        is_serious, detailed_reason = TranslationValidator.is_serious_failure(
            element.full_text, reason, self.large_text_threshold, from_cache)
        with self.failed_tasks_lock:
            failed_task = FailedTask(
                original_text=element.full_text,
                original_index=original_index,
                element=element,
                failure_reason=FailureReason.API_ERROR if is_serious else FailureReason.NOT_TRANSLATED,
                error_message=detailed_reason,
                is_serious=is_serious
            )
            self.failed_tasks.append(failed_task)
            if is_serious:
                self.stats['serious_failures'] += 1
                logger.debug(f"严重失败: {detailed_reason[:50]}...")
            else:
                self.stats['minor_issues'] += 1
                logger.debug(f"轻微问题: {detailed_reason[:50]}...")

    def _add_batch_failure(self, batch: List[PPTElement], batch_indices: List[int], reason: str):
        """添加批次级失败"""
        with self.failed_tasks_lock:
            for element, original_index in zip(batch, batch_indices):
                failed_task = FailedTask(
                    original_text=element.full_text,
                    original_index=original_index,
                    element=element,
                    failure_reason=FailureReason.BATCH_FAILURE,
                    error_message=reason,
                    is_serious=True
                )
                self.failed_tasks.append(failed_task)
                self.stats['serious_failures'] += 1
            logger.warning(f"批次失败: {len(batch)} 个任务, 原因: {reason}")

    def _create_retry_batches(self, failed_tasks: List[FailedTask], retry_count: int) -> List[List[FailedTask]]:
        """创建重试批次，针对PPT优化"""
        serious_failed_tasks = [task for task in failed_tasks if task.is_serious]
        if not serious_failed_tasks:
            return []
        if retry_count < len(self.retry_batch_sizes):
            max_batch_size = self.retry_batch_sizes[retry_count]
        else:
            max_batch_size = 1
        logger.info(f"第{retry_count + 1}次重试，严重失败任务: {len(serious_failed_tasks)}, 批次大小: {max_batch_size}")
        chunks = []
        current_chunk = []
        for task in serious_failed_tasks:
            if len(current_chunk) >= max_batch_size:
                chunks.append(current_chunk)
                current_chunk = [task]
            else:
                current_chunk.append(task)
        if current_chunk:
            chunks.append(current_chunk)
        return chunks

    def _retry_failed_tasks(self, target_lang: str, source_lang: Optional[str]):
        """重试失败任务，支持并发重试"""
        retry_count = 0
        while retry_count < self.max_retries:
            with self.failed_tasks_lock:
                serious_failed_tasks = [task for task in self.failed_tasks 
                                      if task.is_serious and task.retry_count <= retry_count]
            if not serious_failed_tasks:
                logger.info(f"第 {retry_count + 1} 次重试检查：没有严重失败任务需要重试")
                break
            logger.info(f"第 {retry_count + 1} 次重试，处理 {len(serious_failed_tasks)} 个严重失败任务")
            # 清理失败任务缓存
            logger.info(f"清理 {len(serious_failed_tasks)} 个失败任务的缓存")
            for task in serious_failed_tasks:
                self._clear_failed_task_cache(task, target_lang, source_lang)
            # 清空当前重试轮次的失败任务
            with self.failed_tasks_lock:
                self.failed_tasks = [task for task in self.failed_tasks 
                                   if not task.is_serious or task.retry_count > retry_count]
            # 重试延迟
            if retry_count < len(self.retry_delays):
                delay = self.retry_delays[retry_count]
                logger.info(f"重试前等待 {delay} 秒...")
                time.sleep(delay)
            # 创建重试批次
            retry_batches = self._create_retry_batches(serious_failed_tasks, retry_count)
            if not retry_batches:
                logger.info(f"第 {retry_count + 1} 次重试：没有批次需要处理")
                break
            # 判断是否使用并发重试
            use_concurrent = self._should_use_concurrent_retry(retry_batches)
            actual_workers = min(self.retry_max_workers, len(retry_batches)) if use_concurrent else 1
            self.stats['retry_workers_used'] = actual_workers
            if use_concurrent:
                logger.info(f"启用并发重试，worker数量: {actual_workers}")
                self.stats['concurrent_retry_batches'] = len(retry_batches)
                self._retry_batches_concurrent(retry_batches, target_lang, source_lang, retry_count, actual_workers)
            else:
                logger.info("使用串行重试")
                self._retry_batches_sequential(retry_batches, target_lang, source_lang, retry_count)
            self.stats['retry_attempts'] += 1
            retry_count += 1
        # 统计最终失败
        with self.failed_tasks_lock:
            final_serious_failures = [task for task in self.failed_tasks if task.is_serious]
            self.stats['final_failures'] = len(final_serious_failures)
            if final_serious_failures:
                logger.warning(f"常规重试后仍有 {len(final_serious_failures)} 个严重失败任务")

    def _should_use_concurrent_retry(self, retry_batches):
        """判断是否使用并发重试"""
        return len(retry_batches) >= 2  # PPT相对简单，2个批次就可以并发

    def _retry_batches_concurrent(self, retry_batches: List[List[FailedTask]], target_lang: str, 
                                source_lang: Optional[str], retry_count: int, max_workers: int):
        """并发重试批次"""
        try:
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                with tqdm(total=sum(len(batch) for batch in retry_batches), 
                         desc=f"第{retry_count + 1}次重试(并发)", unit="任务") as pbar:
                    future_to_batch = {}
                    for retry_batch_idx, retry_batch in enumerate(retry_batches):
                        future = executor.submit(
                            self._process_retry_batch, retry_batch, target_lang, source_lang, 
                            retry_batch_idx, retry_count
                        )
                        future_to_batch[future] = retry_batch
                    for future in as_completed(future_to_batch):
                        retry_batch = future_to_batch[future]
                        try:
                            future.result(timeout=self.translation_timeout + 20)
                            pbar.update(len(retry_batch))
                        except Exception as e:
                            logger.error(f"并发重试批次处理异常: {e}")
                            with self.failed_tasks_lock:
                                for task in retry_batch:
                                    new_task = FailedTask(
                                        original_text=task.original_text,
                                        original_index=task.original_index,
                                        element=task.element,
                                        failure_reason=FailureReason.BATCH_FAILURE,
                                        error_message=f"并发执行异常: {str(e)}",
                                        is_serious=True
                                    )
                                    new_task.retry_count = retry_count + 1
                                    self.failed_tasks.append(new_task)
                            pbar.update(len(retry_batch))
        except Exception as e:
            logger.error(f"并发重试执行异常: {e}")
            logger.info("回退到串行重试")
            self._retry_batches_sequential(retry_batches, target_lang, source_lang, retry_count)

    def _retry_batches_sequential(self, retry_batches: List[List[FailedTask]], target_lang: str, 
                                source_lang: Optional[str], retry_count: int):
        """串行重试批次"""
        with tqdm(total=sum(len(batch) for batch in retry_batches), 
                 desc=f"第{retry_count + 1}次重试", unit="任务") as pbar:
            for retry_batch_idx, retry_batch in enumerate(retry_batches):
                self._process_retry_batch(retry_batch, target_lang, source_lang, retry_batch_idx, retry_count)
                pbar.update(len(retry_batch))

    def _process_retry_batch(self, retry_batch: List[FailedTask], target_lang: str, 
                           source_lang: Optional[str], batch_idx: int, retry_count: int):
        """处理单个重试批次"""
        logger.info(f"开始重试批次 {batch_idx + 1}，任务数: {len(retry_batch)}")
        batch_elements = [task.element for task in retry_batch]
        batch_indices = [task.original_index for task in retry_batch]
        success, retry_results, cache_flags = self._translate_batch_with_timeout(
            batch_elements, batch_indices, target_lang, source_lang, -1)
        logger.info(f"重试批次 {batch_idx + 1} 完成，成功: {success}, 结果数: {len(retry_results) if retry_results else 0}")
        if not success:
            batch_reason = retry_results[0] if retry_results else "重试批次失败"
            logger.warning(f"重试批次失败: {batch_reason}")
            with self.failed_tasks_lock:
                for task in retry_batch:
                    new_task = FailedTask(
                        original_text=task.original_text,
                        original_index=task.original_index,
                        element=task.element,
                        failure_reason=FailureReason.BATCH_FAILURE,
                        error_message=batch_reason,
                        is_serious=True
                    )
                    new_task.retry_count = retry_count + 1
                    self.failed_tasks.append(new_task)
        else:
            for i, task in enumerate(retry_batch):
                if i < len(retry_results):
                    result = retry_results[i]
                    from_cache = cache_flags[i] if i < len(cache_flags) else False
                    is_serious, reason = TranslationValidator.is_serious_failure(
                        task.original_text, result, self.large_text_threshold, from_cache)
                    if not is_serious and not TranslationValidator.is_error_message(result):
                        self._apply_translation_to_ppt_element(task.element, result)
                        logger.debug(f"重试成功: {task.original_text[:50]}... -> {result[:50]}...")
                    else:
                        logger.debug(f"重试失败，保持原文: {task.original_text[:50]}... 原因: {reason}")
                        with self.failed_tasks_lock:
                            new_task = FailedTask(
                                original_text=task.original_text,
                                original_index=task.original_index,
                                element=task.element,
                                failure_reason=FailureReason.API_ERROR,
                                error_message=reason,
                                is_serious=True
                            )
                            new_task.retry_count = retry_count + 1
                            self.failed_tasks.append(new_task)
                else:
                    logger.warning(f"重试结果不足: 任务 {i}, 结果数量 {len(retry_results)}")
                    with self.failed_tasks_lock:
                        new_task = FailedTask(
                            original_text=task.original_text,
                            original_index=task.original_index,
                            element=task.element,
                            failure_reason=FailureReason.PARSE_ERROR,
                            error_message="重试结果缺失",
                            is_serious=True
                        )
                        new_task.retry_count = retry_count + 1
                        self.failed_tasks.append(new_task)

    def _final_retry_remaining_tasks(self, target_lang: str, source_lang: Optional[str]):
        """最终处理剩余失败任务"""
        with self.failed_tasks_lock:
            remaining_tasks = [task for task in self.failed_tasks if task.is_serious]
        if not remaining_tasks:
            logger.info("没有剩余失败任务需要最终处理")
            return
        logger.info(f"最终处理 {len(remaining_tasks)} 个剩余失败任务")
        success_count = 0
        with tqdm(total=len(remaining_tasks), desc="最终挽救", unit="任务") as pbar:
            for task in remaining_tasks:
                try:
                    self._clear_failed_task_cache(task, target_lang, source_lang)
                    time.sleep(0.8)  # PPT处理间隔稍短
                    logger.debug(f"最终处理任务: {task.original_text[:50]}...")
                    result = self.translator.translate(
                        text=task.original_text,
                        target_lang=target_lang,
                        source_lang=source_lang
                    )
                    self.stats['api_calls'] += 1
                    if (result and not TranslationValidator.is_error_message(result) and 
                        result.strip() != task.original_text.strip()):
                        is_serious, reason = TranslationValidator.is_serious_failure(
                            task.original_text, result, self.large_text_threshold, False)
                        if not is_serious:
                            self._apply_translation_to_ppt_element(task.element, result)
                            success_count += 1
                            self.stats['final_rescues'] += 1
                            logger.info(f"最终处理成功: {task.original_text[:50]}... -> {result[:50]}...")
                        else:
                            logger.debug(f"最终处理结果仍不合格: {reason}")
                    else:
                        logger.debug(f"最终处理仍失败: {task.original_text[:50]}...")
                except Exception as e:
                    logger.error(f"最终处理异常: {e}")
                pbar.update(1)
        if success_count > 0:
            logger.info(f"最终挽救成功 {success_count}/{len(remaining_tasks)} 个任务")
        else:
            logger.warning(f"最终挽救失败，{len(remaining_tasks)} 个任务保持原文")

    def _collect_ppt_elements(self, presentation: Presentation) -> List[PPTElement]:
        """收集PPT中的可翻译元素"""
        elements = []
        try:
            for slide_idx, slide in enumerate(presentation.slides):
                self.stats['total_slides'] += 1
                # 处理幻灯片中的形状
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_elements = self._extract_text_from_shape(shape, slide_idx, shape_idx)
                    elements.extend(shape_elements)
                # 处理备注页（如果启用）
                if self.translate_notes and hasattr(slide, 'notes_slide') and slide.notes_slide:
                    try:
                        notes_text_frame = slide.notes_slide.notes_text_frame
                        if notes_text_frame and notes_text_frame.text.strip():
                            # 提取备注页格式信息
                            notes_format = ParagraphFormatInfo()
                            if notes_text_frame.paragraphs:
                                notes_format = _extract_paragraph_format_info_enhanced(notes_text_frame.paragraphs[0])
                            
                            elements.append(PPTElement(
                                full_text=notes_text_frame.text,
                                slide_index=slide_idx,
                                shape=None,
                                shape_type='notes',
                                shape_info={'notes': True},
                                text_container=notes_text_frame,
                                location=f"slide_{slide_idx}_notes",
                                format_info=notes_format,
                                element_type='notes'
                            ))
                            self.stats['notes_pages'] += 1
                    except Exception as e:
                        logger.debug(f"处理备注页异常: {e}")
        except Exception as e:
            logger.error(f"Error collecting PPT elements: {e}")
        return elements

    def _extract_text_from_shape(self, shape: BaseShape, slide_idx: int, shape_idx: int) -> List[PPTElement]:
        """从形状中提取文本元素（增强格式保持版）"""
        elements = []
        try:
            # 检查形状类型
            shape_type = self._get_shape_type(shape)
            if shape_type == 'text' and hasattr(shape, 'text_frame'):
                # 处理文本框和占位符
                text_frame = shape.text_frame
                if text_frame and text_frame.text.strip():
                    # 提取形状信息
                    shape_info = self._extract_shape_info(shape)
                    # 逐段落处理以保持格式
                    for para_idx, paragraph in enumerate(text_frame.paragraphs):
                        if paragraph.text.strip():
                            # 提取增强的格式信息
                            format_info = _extract_paragraph_format_info_enhanced(paragraph)
                            elements.append(PPTElement(
                                full_text=paragraph.text,
                                slide_index=slide_idx,
                                shape=shape,
                                shape_type=shape_type,
                                shape_info=shape_info,
                                text_container=text_frame,
                                paragraph_index=para_idx,
                                location=f"slide_{slide_idx}_shape_{shape_idx}_para_{para_idx}",
                                format_info=format_info,
                                element_type='text_paragraph'
                            ))
                    # 更新统计
                    if hasattr(shape, 'placeholder_format'):
                        self.stats['placeholders'] += 1
                    else:
                        self.stats['text_boxes'] += 1
            elif shape_type == 'table' and hasattr(shape, 'table'):
                # 处理表格
                table = shape.table
                self.stats['tables'] += 1
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text_frame and cell.text_frame.text.strip():
                            cell_info = {
                                'row': row_idx,
                                'col': col_idx,
                                'table_size': (len(table.rows), len(table.columns))
                            }
                            # 提取表格单元格的格式信息
                            cell_format = ParagraphFormatInfo()
                            if cell.text_frame.paragraphs:
                                cell_format = _extract_paragraph_format_info_enhanced(cell.text_frame.paragraphs[0])
                            
                            elements.append(PPTElement(
                                full_text=cell.text_frame.text,
                                slide_index=slide_idx,
                                shape=shape,
                                shape_type='table_cell',
                                shape_info=self._extract_shape_info(shape),
                                text_container=cell.text_frame,
                                cell_info=cell_info,
                                location=f"slide_{slide_idx}_table_{shape_idx}_cell_{row_idx}_{col_idx}",
                                format_info=cell_format,
                                element_type='table_cell'
                            ))
            elif shape_type == 'chart':
                # 处理图表（如果有文本）
                self.stats['charts'] += 1
                chart_elements = self._extract_chart_text(shape, slide_idx, shape_idx)
                elements.extend(chart_elements)
            elif shape_type == 'smartart':
                # 处理SmartArt图形
                self.stats['smartart'] += 1
                # SmartArt的文本提取比较复杂，这里简化处理
                # 实际实现中可能需要更复杂的逻辑
                pass
        except Exception as e:
            logger.debug(f"提取形状文本异常: {e}")
        return elements

    def _get_shape_type(self, shape: BaseShape) -> str:
        """确定形状类型"""
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                return 'text'
            elif hasattr(shape, 'table') and shape.table:
                return 'table'
            elif hasattr(shape, 'chart') and shape.chart:
                return 'chart'
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                return 'group'
            elif 'smartart' in str(type(shape)).lower():
                return 'smartart'
            else:
                return 'other'
        except:
            return 'unknown'

    def _extract_shape_info(self, shape: BaseShape) -> Dict[str, Any]:
        """提取形状信息"""
        info = {}
        try:
            if hasattr(shape, 'left'):
                info['left'] = shape.left
            if hasattr(shape, 'top'):
                info['top'] = shape.top
            if hasattr(shape, 'width'):
                info['width'] = shape.width
            if hasattr(shape, 'height'):
                info['height'] = shape.height
            if hasattr(shape, 'rotation'):
                info['rotation'] = shape.rotation
            if hasattr(shape, 'name'):
                info['name'] = shape.name
        except:
            pass
        return info

    def _extract_chart_text(self, shape: BaseShape, slide_idx: int, shape_idx: int) -> List[PPTElement]:
        """提取图表中的文本"""
        elements = []
        try:
            # 图表文本提取比较复杂，这里提供基础框架
            # 实际实现需要根据具体的图表类型来处理
            if hasattr(shape, 'chart'):
                chart = shape.chart
                # 这里可以提取图表标题、轴标签、图例等文本
                # 由于python-pptx对图表的支持有限，这里简化处理
                pass
        except Exception as e:
            logger.debug(f"提取图表文本异常: {e}")
        return elements

    def _apply_translation_to_ppt_element(self, element: PPTElement, translated_text: str):
        """将翻译结果应用到PPT元素（格式保持增强版）"""
        if not translated_text or translated_text == element.full_text:
            return
        
        try:
            if element.element_type == 'text_paragraph':
                # 处理文本段落 - 关键改进点
                if element.paragraph_index is not None and element.text_container:
                    text_frame = element.text_container
                    if element.paragraph_index < len(text_frame.paragraphs):
                        paragraph = text_frame.paragraphs[element.paragraph_index]
                        
                        if self.preserve_formatting and element.format_info:
                            # 使用格式保持模式
                            try:
                                _rebuild_paragraph_with_format(paragraph, translated_text, element.format_info)
                                self.stats['format_preserved'] += 1
                                logger.debug(f"格式保持成功: {element.full_text[:30]}... -> {translated_text[:30]}...")
                            except Exception as e:
                                logger.debug(f"格式保持失败，使用基础模式: {e}")
                                self._apply_basic_translation(paragraph, translated_text)
                                self.stats['format_fallback'] += 1
                        else:
                            # 使用基础模式
                            self._apply_basic_translation(paragraph, translated_text)
                            
            elif element.element_type == 'table_cell':
                # 处理表格单元格
                if element.text_container:
                    text_frame = element.text_container
                    if self.preserve_formatting and element.format_info:
                        # 表格单元格格式保持
                        try:
                            text_frame.clear()
                            paragraph = text_frame.paragraphs[0]
                            _rebuild_paragraph_with_format(paragraph, translated_text, element.format_info)
                            self.stats['format_preserved'] += 1
                        except Exception as e:
                            logger.debug(f"表格单元格格式保持失败: {e}")
                            text_frame.clear()
                            paragraph = text_frame.paragraphs[0]
                            run = paragraph.add_run()
                            run.text = translated_text
                            self.stats['format_fallback'] += 1
                    else:
                        # 基础表格单元格处理
                        text_frame.clear()
                        paragraph = text_frame.paragraphs[0]
                        run = paragraph.add_run()
                        run.text = translated_text
                        
            elif element.element_type == 'notes':
                # 处理备注页
                if element.text_container:
                    if self.preserve_formatting and element.format_info:
                        try:
                            element.text_container.clear()
                            paragraph = element.text_container.paragraphs[0]
                            _rebuild_paragraph_with_format(paragraph, translated_text, element.format_info)
                            self.stats['format_preserved'] += 1
                        except Exception as e:
                            logger.debug(f"备注页格式保持失败: {e}")
                            element.text_container.clear()
                            paragraph = element.text_container.paragraphs[0]
                            run = paragraph.add_run()
                            run.text = translated_text
                            self.stats['format_fallback'] += 1
                    else:
                        element.text_container.clear()
                        paragraph = element.text_container.paragraphs[0]
                        run = paragraph.add_run()
                        run.text = translated_text
                        
            self.stats['translated_elements'] += 1
            
        except Exception as e:
            logger.warning(f"Error applying translation to PPT element: {e}")
            # 失败时的最后尝试
            try:
                self._apply_fallback_translation(element, translated_text)
            except Exception as fallback_e:
                logger.error(f"Fallback translation also failed: {fallback_e}")

    def _apply_basic_translation(self, paragraph, translated_text: str):
        """应用基础翻译（无格式保持）"""
        try:
            # 清空段落并添加翻译文本
            paragraph.clear()
            run = paragraph.add_run()
            run.text = translated_text
        except Exception as e:
            logger.debug(f"基础翻译应用失败: {e}")
            raise

    def _apply_fallback_translation(self, element: PPTElement, translated_text: str):
        """最后的翻译应用方法"""
        try:
            if element.element_type == 'text_paragraph' and element.text_container:
                text_frame = element.text_container
                if element.paragraph_index is not None and element.paragraph_index < len(text_frame.paragraphs):
                    paragraph = text_frame.paragraphs[element.paragraph_index]
                    paragraph.clear()
                    run = paragraph.add_run()
                    run.text = translated_text
            elif element.text_container and hasattr(element.text_container, 'clear'):
                element.text_container.clear()
                if hasattr(element.text_container, 'paragraphs') and element.text_container.paragraphs:
                    paragraph = element.text_container.paragraphs[0]
                    run = paragraph.add_run()
                    run.text = translated_text
        except Exception as e:
            logger.debug(f"Fallback translation failed: {e}")


# pptx_full_translator.py (第三批 - 格式保持增强版)

    def _get_cache_key(self, text: str, target_lang: str, source_lang: Optional[str], 
                      prompt_config: Optional[Dict[str, Any]] = None) -> str:
        """生成缓存键，包含prompt配置信息"""
        try:
            config = prompt_config or self.effective_prompt_config
            if self._cached_config_hash is None and config is not None:
                try:
                    mode = config.get('mode', '') if config else ''
                    template = config.get('prompt_template', '') if config else ''
                    custom_system = ''
                    custom_prompt = config.get('custom_prompt')
                    if custom_prompt and isinstance(custom_prompt, dict):
                        custom_system = custom_prompt.get('system', '')[:50]
                    key_config_str = f"{mode}_{template}_{custom_system}"
                    self._cached_config_hash = hashlib.md5(key_config_str.encode('utf-8')).hexdigest()[:8]
                except (AttributeError, TypeError, KeyError) as e:
                    logger.warning(f"Failed to generate prompt hash, using empty: {e}")
                    self._cached_config_hash = ""
            prompt_hash = self._cached_config_hash or ""
            text_hash = hashlib.md5(text.encode('utf-8')).hexdigest()
            return f"{text_hash}_{target_lang}_{source_lang or 'auto'}_{prompt_hash}"
        except Exception as e:
            logger.warning(f"Error generating cache key: {e}")
            text_hash = hashlib.md5(text.encode('utf-8')).hexdigest()
            return f"{text_hash}_{target_lang}_{source_lang or 'auto'}"

    def _get_enhanced_system_prompt(self, target_lang: str, source_lang: Optional[str]) -> str:
        """获取增强的系统提示，适配PPT翻译"""
        try:
            if not self.effective_prompt_config or self.effective_prompt_config.get('mode') == 'none':
                return self._build_default_ppt_system_prompt(target_lang, source_lang)
            mode = self.effective_prompt_config.get('mode', 'none')
            if mode == 'custom' and self.effective_prompt_config.get('custom_prompt'):
                return self._build_custom_ppt_system_prompt(target_lang, source_lang)
            elif mode == 'professional':
                return self._build_professional_ppt_system_prompt(target_lang, source_lang)
            elif mode == 'general':
                return self._build_general_ppt_system_prompt(target_lang, source_lang)
            else:
                return self._build_simple_ppt_system_prompt(target_lang, source_lang)
        except Exception as e:
            logger.warning(f"Error generating enhanced system prompt: {e}")
            return self._build_default_ppt_system_prompt(target_lang, source_lang)

    def _build_custom_ppt_system_prompt(self, target_lang: str, source_lang: Optional[str]) -> str:
        """构建自定义PPT系统提示"""
        try:
            custom_prompt = self.effective_prompt_config.get('custom_prompt', {})
            system_content = custom_prompt.get('system', '')
            if not system_content:
                logger.warning("Custom prompt system content is empty, falling back to default")
                return self._build_default_ppt_system_prompt(target_lang, source_lang)
            # 为PPT批量翻译添加必要说明
            if "numbered line" not in system_content.lower():
                system_content += f"""

ADVANCED PPT BATCH PROCESSING WITH FORMAT PRESERVATION:
- Each input line is numbered [1], [2], etc. representing different presentation elements
- Translate each numbered line individually while maintaining presentation context
- Preserve slide structure and formatting context
- Keep the exact same number of lines as input
- Output only the translated content without including the original text, one per line
- Do not include line numbers in output
- Maintain presentation style and visual impact
- Preserve bullet points and list formatting
- Consider that formatting will be automatically preserved by the system
- Do not include any extra comments in your output
- Ensure translations sound natural in {target_lang} for presentation context"""
            system_content = self._add_ppt_enhancement_rules(system_content)
            logger.info("Using custom prompt for advanced PPT translation with format preservation")
            return system_content
        except Exception as e:
            logger.warning(f"Error building custom PPT system prompt: {e}")
            return self._build_default_ppt_system_prompt(target_lang, source_lang)

    def _build_professional_ppt_system_prompt(self, target_lang: str, source_lang: Optional[str]) -> str:
        """构建专业PPT模板系统提示"""
        try:
            domain = self.effective_prompt_config.get('prompt_template', 'business')
            logger.info(f"Using professional template for advanced PPT: {domain}")
            # PPT专业领域的系统提示
            professional_prompts = {
                'business': f"""You are an expert presentation translator specializing in business presentations.
Translate from {source_lang or 'auto-detected language'} to {target_lang}.
Maintain professional presentation style, use impactful business language, and preserve visual hierarchy.
Focus on:
- Business presentation terminology and corporate language
- Clear, concise, and impactful messaging
- Professional communication style suitable for executives
- Maintaining slide flow and narrative structure
- Preserving bullet points and key messaging
- Format preservation is handled automatically""",
                'academic': f"""You are an academic presentation translator with expertise in scholarly presentations.
Translate from {source_lang or 'auto-detected language'} to {target_lang}.
Maintain academic rigor, preserve research terminology, and keep citation formats intact.
Pay attention to:
- Academic presentation style and formal tone
- Research methodology and findings presentation
- Technical and disciplinary terminology
- Conference presentation standards
- Preserving data visualization context
- Format preservation is handled automatically""",
                'technical': f"""You are a technical presentation translator specializing in technical demos and documentation.
Translate from {source_lang or 'auto-detected language'} to {target_lang}.
Preserve technical accuracy, maintain procedural clarity, and use industry-standard terminology.
Focus on:
- Technical demonstration language
- Step-by-step procedure clarity
- Software and hardware terminology
- Product feature descriptions
- Maintaining technical diagram context
- Format preservation is handled automatically""",
                'sales': f"""You are a sales presentation translator with expertise in persuasive presentations.
Translate from {source_lang or 'auto-detected language'} to {target_lang}.
Maintain persuasive impact, preserve sales messaging, and adapt cultural context appropriately.
Consider:
- Sales pitch terminology and persuasive language
- Value proposition clarity
- Customer benefit messaging
- Call-to-action effectiveness
- Maintaining emotional impact
- Format preservation is handled automatically""",
                'education': f"""You are an educational presentation translator focusing on learning materials.
Translate from {source_lang or 'auto-detected language'} to {target_lang}.
Maintain educational clarity, preserve learning objectives, and ensure student comprehension.
Pay attention to:
- Educational terminology and concepts
- Learning objective clarity
- Student-friendly language
- Instructional design principles
- Maintaining pedagogical structure
- Format preservation is handled automatically"""
            }
            system_content = professional_prompts.get(domain, professional_prompts['business'])
            # 添加PPT批处理规则
            system_content += f"""

PROFESSIONAL PPT BATCH PROCESSING WITH FORMAT PRESERVATION:
1. Process each numbered line [1], [2], etc. individually for different slide elements
2. Consider the presentation context and slide flow for coherent translation
3. Maintain professional presentation consistency throughout
4. Keep the exact same number of lines as input
5. Output only the translated content, one per line
6. Do not include line numbers in output
7. Preserve professional presentation formatting and impact
8. Do not include any extra comments in your output
9. Ensure translations maintain presentation effectiveness in {target_lang}
10. Font styles, colors, and formatting will be preserved automatically"""
            system_content = self._add_ppt_enhancement_rules(system_content)
            return system_content
        except Exception as e:
            logger.warning(f"Error building professional PPT system prompt: {e}")
            return self._build_default_ppt_system_prompt(target_lang, source_lang)

    def _build_general_ppt_system_prompt(self, target_lang: str, source_lang: Optional[str]) -> str:
        """构建通用PPT增强系统提示"""
        try:
            system_content = f"""You are a professional presentation translator with expertise in slide translation.
Translate from {source_lang or 'auto-detected language'} to {target_lang}.
Provide accurate, natural translations while preserving presentation impact and visual hierarchy.
Maintain slide coherence and consistency throughout the presentation.

GENERAL PPT BATCH PROCESSING RULES WITH FORMAT PRESERVATION:
1. Translate each numbered line [1], [2], etc. representing different slide elements
2. Consider context from surrounding elements for coherent presentation flow
3. Keep the exact same number of lines as input
4. Preserve all formatting, bullet points, and visual structure (handled automatically)
5. Output only the translated content, one per line
6. Do not include line numbers in output
7. Maintain consistency in terminology throughout the presentation
8. Preserve presentation impact and messaging effectiveness
9. Do not include any extra comments in your output
10. Ensure translations sound natural in {target_lang} for presentation context
11. Font styles, colors, and formatting will be preserved by the system"""
            system_content = self._add_ppt_enhancement_rules(system_content)
            return system_content
        except Exception as e:
            logger.warning(f"Error building general PPT system prompt: {e}")
            return self._build_default_ppt_system_prompt(target_lang, source_lang)

    def _build_simple_ppt_system_prompt(self, target_lang: str, source_lang: Optional[str]) -> str:
        """构建简单PPT系统提示"""
        return f"""Translate presentation content from {source_lang or 'auto-detected language'} to {target_lang}.
Process each numbered line [1], [2], etc. and return the same number of translated lines.
Maintain presentation style and impact. Do not include line numbers in output.
Formatting will be preserved automatically."""

    def _build_default_ppt_system_prompt(self, target_lang: str, source_lang: Optional[str]) -> str:
        """构建默认PPT系统提示"""
        format_note = "Formatting (fonts, colors, styles) will be preserved automatically by the system." if self.preserve_formatting else ""
        
        system_content = f"""You are a professional presentation translator. Translate from {source_lang or 'auto-detected language'} to {target_lang}.

Rules for presentation translation:
1. Translate each numbered line individually, but consider the full presentation context
2. Keep the exact same number of lines as the original
3. Preserve all formatting, bullet points, and presentation structure
4. For slide titles, maintain impact and clarity
5. For bullet points, keep the list structure and visual hierarchy
6. Output only the translated lines, one per line, in the same order
7. Do not include the original line numbers or any extra comments in your output
8. Do not translate company names, product names, or proper nouns—keep them exactly as in the original
9. Maintain presentation flow and messaging effectiveness
10. Return only the translated content without including the original text
{format_note}"""
        return system_content

    def _add_ppt_enhancement_rules(self, system_content: str) -> str:
        """添加PPT特有的增强规则到系统提示"""
        try:
            if not self.effective_prompt_config or self.effective_prompt_config.get('mode') == 'none':
                return system_content
            enhancements = []
            # 保留术语
            preserve_terms = self.effective_prompt_config.get('preserve_terms')
            if preserve_terms and isinstance(preserve_terms, list) and preserve_terms:
                terms = ', '.join(str(term) for term in preserve_terms if term)
                if terms:
                    enhancements.append(f"PRESERVE THESE TERMS EXACTLY: {terms}")
            # 术语表
            glossary = self.effective_prompt_config.get('glossary')
            if glossary and isinstance(glossary, dict) and glossary:
                try:
                    glossary_text = '; '.join([f"{k}: {v}" for k, v in glossary.items() if k and v])
                    if glossary_text:
                        enhancements.append(f"USE THIS GLOSSARY: {glossary_text}")
                except Exception as e:
                    logger.warning(f"Error processing glossary: {e}")
            # 额外上下文
            additional_context = self.effective_prompt_config.get('additional_context')
            if additional_context and str(additional_context).strip():
                enhancements.append(f"PRESENTATION CONTEXT: {str(additional_context).strip()}")
            if enhancements:
                enhancement_text = "\n\nADDITIONAL PRESENTATION REQUIREMENTS:\n" + "\n".join(f"• {rule}" for rule in enhancements)
                system_content += enhancement_text
            return system_content
        except Exception as e:
            logger.warning(f"Error adding PPT enhancement rules: {e}")
            return system_content

    def _create_ordered_batches(self, elements: List[PPTElement]) -> List[Tuple[List[PPTElement], List[int]]]:
        """创建批次，保持幻灯片顺序"""
        batches = []
        current_batch = []
        current_indices = []
        current_chars = 0
        for idx, element in enumerate(elements):
            text_len = len(element.full_text)
            if (len(current_batch) >= self.batch_size or 
                current_chars + text_len > self.max_chars) and current_batch:
                batches.append((current_batch, current_indices))
                current_batch = [element]
                current_indices = [idx]
                current_chars = text_len
            else:
                current_batch.append(element)
                current_indices.append(idx)
                current_chars += text_len
        if current_batch:
            batches.append((current_batch, current_indices))
        return batches

    def _translate_batch_with_timeout(self, batch: List[PPTElement], batch_indices: List[int],
                                    target_lang: str, source_lang: Optional[str], 
                                    batch_num: int) -> Tuple[bool, List[str], List[bool]]:
        """带超时的批次翻译"""
        try:
            with ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(self._translate_batch, batch, target_lang, source_lang, batch_num)
                try:
                    batch_results, from_cache_flags = future.result(timeout=self.translation_timeout)
                    return True, batch_results, from_cache_flags
                except TimeoutError:
                    logger.warning(f"批次 {batch_num} 翻译超时 ({self.translation_timeout}秒)")
                    failure_reasons = [f"翻译超时 ({self.translation_timeout}秒)"] * len(batch)
                    cache_flags = [False] * len(batch)
                    return False, failure_reasons, cache_flags
        except Exception as e:
            logger.error(f"批次 {batch_num} 翻译异常: {e}")
            failure_reasons = [f"翻译异常: {str(e)}"] * len(batch)
            cache_flags = [False] * len(batch)
            return False, failure_reasons, cache_flags

    def _translate_batch(self, batch: List[PPTElement], target_lang: str, 
                        source_lang: Optional[str], batch_num: int) -> Tuple[List[str], List[bool]]:
        """翻译单个批次"""
        texts = [element.full_text for element in batch]
        # 检查缓存
        results = []
        uncached_texts = []
        uncached_indices = []
        local_prompt_config = self._get_config_safe()
        for i, text in enumerate(texts):
            cache_key = self._get_cache_key(text, target_lang, source_lang, local_prompt_config)
            cached = self.cache.get(cache_key)
            if cached:
                results.append(cached)
                self.stats['cache_savings'] += 1
            else:
                uncached_texts.append(text)
                uncached_indices.append(i)
                results.append("")
        # 翻译未缓存的文本
        if uncached_texts:
            try:
                system_prompt = self._get_enhanced_system_prompt(target_lang, source_lang)
                # 使用行号编码
                numbered_texts = [f"[{i+1}] {text}" for i, text in enumerate(uncached_texts)]
                user_message = "\n".join(numbered_texts)
                # 调用翻译器
                try:
                    translated_result = self.translator.translate(
                        text=user_message,
                        target_lang=target_lang,
                        source_lang=source_lang,
                        prompt_config=local_prompt_config,
                        config_merge_mode='merge'
                    )
                except (TypeError, AttributeError):
                    try:
                        messages = [
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_message}
                        ]
                        translated_result = self.translator.translate(
                            messages=messages,
                            target_lang=target_lang,
                            source_lang=source_lang
                        )
                    except:
                        full_prompt = f"{system_prompt}\n\nPresentation content to translate:\n{user_message}"
                        translated_result = self.translator.translate(full_prompt, target_lang, source_lang)
                self.stats['api_calls'] += 1
                # 解析结果
                translated_parts = self._extract_numbered_translations(translated_result, len(uncached_texts))
                # 更新缓存和结果
                for i, (original_text, idx) in enumerate(zip(uncached_texts, uncached_indices)):
                    if i < len(translated_parts) and translated_parts[i]:
                        translation = translated_parts[i].strip()
                        cache_key = self._get_cache_key(original_text, target_lang, source_lang, local_prompt_config)
                        self.cache.put(cache_key, translation)
                        results[idx] = translation
                    else:
                        results[idx] = original_text
            except Exception as e:
                logger.error(f"批次 {batch_num} 翻译失败: {e}")
                for idx in uncached_indices:
                    results[idx] = texts[idx]
        # 返回结果和缓存标记
        from_cache_flags = [bool(results[i] and i not in uncached_indices) for i in range(len(results))]
        return results, from_cache_flags

    def _extract_numbered_translations(self, response: str, expected_count: int) -> List[str]:
        """提取编号翻译结果"""
        translations = [""] * expected_count
        try:
            lines = response.strip().split('\n')
            current_line_num = 0
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                # 检查行号标记
                number_match = NUMBERED_LINE_PATTERN.match(line)
                if number_match:
                    line_num = int(number_match.group(1)) - 1
                    content = number_match.group(2).strip()
                    if 0 <= line_num < expected_count and content:
                        translations[line_num] = content
                else:
                    # 无行号标记，按顺序分配
                    if current_line_num < expected_count and line:
                        translations[current_line_num] = line
                        current_line_num += 1
            # 填补空缺
            if any(not t for t in translations):
                non_empty_lines = [line.strip() for line in lines if line.strip()]
                for i, line in enumerate(non_empty_lines):
                    if i < expected_count and not translations[i]:
                        clean_line = re.sub(r'^\[\d+\]\s*', '', line)
                        if clean_line:
                            translations[i] = clean_line
        except Exception as e:
            logger.warning(f"Error extracting numbered translations: {e}")
        return translations

    def translate_pptx(self, input_filepath: str, output_filepath: str, target_lang: str, source_lang: Optional[str] = None) -> str:
        """主翻译方法（格式保持增强版）"""
        try:
            start_time = time.time()
            self.source_lang = source_lang
            with self._translator_config_context():
                logger.info(f"开始高级PPTX翻译（格式保持增强版）: {os.path.basename(input_filepath)}")
                logger.info(f"目标语言: {target_lang}")
                logger.info(f"格式保持模式: {'启用' if self.preserve_formatting else '禁用'}")
                
                if self.effective_prompt_config and self.effective_prompt_config.get('mode') != 'none':
                    mode = self.effective_prompt_config.get('mode', 'none')
                    logger.info(f"Prompt配置: mode={mode}")
                    if mode == 'professional':
                        logger.info(f"专业领域: {self.effective_prompt_config.get('prompt_template', 'business')}")
                    elif mode == 'custom':
                        logger.info("使用自定义prompt")
                        
                logger.info(f"PPT设置: 翻译备注={self.translate_notes}, 翻译标题={self.translate_slide_titles}")
                logger.info(f"批处理设置: batch_size={self.batch_size}, max_chars={self.max_chars}, workers={self.max_workers}")
                logger.info(f"重试设置: max_retries={self.max_retries}, timeout={self.translation_timeout}s, retry_workers={self.retry_max_workers}")
                
                # 复制原文件
                shutil.copy2(input_filepath, output_filepath)
                presentation = Presentation(output_filepath)
                
                # 收集PPT元素
                elements = self._collect_ppt_elements(presentation)
                if not elements:
                    logger.info("没有需要翻译的内容")
                    print("翻译完成！（无需要翻译的内容）")
                    return output_filepath
                
                self.stats['total_elements'] = len(elements)
                self.stats['total_chars'] = sum(len(e.full_text) for e in elements)
                
                # 创建批次
                batches = self._create_ordered_batches(elements)
                self.stats['total_batches'] = len(batches)
                
                logger.info(f"需要翻译 {len(elements)} 个元素")
                logger.info(f"分为 {len(batches)} 个批次处理")
                
                # 清空失败任务列表
                with self.failed_tasks_lock:
                    self.failed_tasks.clear()
                
                # 执行翻译
                self._translate_all_batches(batches, target_lang, source_lang)
                
                # 智能重试
                if self._should_trigger_retry(len(elements)):
                    logger.info(f"检测到严重失败，开始重试流程")
                    self._retry_failed_tasks(target_lang, source_lang)
                else:
                    with self.failed_tasks_lock:
                        serious_count = sum(1 for task in self.failed_tasks if task.is_serious)
                        minor_count = sum(1 for task in self.failed_tasks if not task.is_serious)
                        if self.failed_tasks:
                            logger.info(f"失败分析: 严重失败 {serious_count} 个, 轻微问题 {minor_count} 个, 未达到重试阈值，跳过重试")
                
                # 最终处理
                self._final_retry_remaining_tasks(target_lang, source_lang)
                
                # 保存文件
                presentation.save(output_filepath)
                self.stats['processing_time'] = time.time() - start_time
                self._print_stats()
                return output_filepath
                
        except Exception as e:
            logger.error(f"高级PPTX翻译失败: {e}")
            raise

    def _translate_all_batches(self, batches: List[Tuple[List[PPTElement], List[int]]], 
                             target_lang: str, source_lang: Optional[str]):
        """翻译所有批次"""
        if len(batches) > 8 and self.max_workers > 1:  # PPT阈值调低
            self._translate_concurrent(batches, target_lang, source_lang)
        else:
            self._translate_sequential(batches, target_lang, source_lang)

    def _translate_sequential(self, batches: List[Tuple[List[PPTElement], List[int]]], 
                            target_lang: str, source_lang: Optional[str]):
        """串行翻译"""
        with tqdm(total=sum(len(batch[0]) for batch in batches), 
                 desc="高级PPTX翻译进度", unit="元素") as pbar:
            for batch_idx, (batch_elements, batch_indices) in enumerate(batches):
                success, translated_texts, from_cache_flags = self._translate_batch_with_timeout(
                    batch_elements, batch_indices, target_lang, source_lang, batch_idx + 1)
                if success:
                    for i, (element, translated_text) in enumerate(zip(batch_elements, translated_texts)):
                        from_cache = from_cache_flags[i] if i < len(from_cache_flags) else False
                        is_serious, reason = TranslationValidator.is_serious_failure(
                            element.full_text, translated_text, self.large_text_threshold, from_cache)
                        if is_serious or (not from_cache and reason != "翻译成功"):
                            self._add_failed_task(element, batch_indices[i], translated_text, from_cache)
                        self._apply_translation_to_ppt_element(element, translated_text)
                else:
                    self._add_batch_failure(batch_elements, batch_indices, translated_texts[0] if translated_texts else "批次处理失败")
                pbar.update(len(batch_elements))

    def _translate_concurrent(self, batches: List[Tuple[List[PPTElement], List[int]]], 
                            target_lang: str, source_lang: Optional[str]):
        """并发翻译"""
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            with tqdm(total=sum(len(batch[0]) for batch in batches), 
                     desc="高级PPTX翻译进度", unit="元素") as pbar:
                future_to_batch = {}
                for batch_idx, (batch_elements, batch_indices) in enumerate(batches):
                    future = executor.submit(
                        self._translate_batch_with_timeout, batch_elements, batch_indices,
                        target_lang, source_lang, batch_idx + 1
                    )
                    future_to_batch[future] = (batch_elements, batch_indices)
                for future in as_completed(future_to_batch):
                    batch_elements, batch_indices = future_to_batch[future]
                    try:
                        success, translated_texts, from_cache_flags = future.result()
                        if success:
                            for i, (element, translated_text) in enumerate(zip(batch_elements, translated_texts)):
                                from_cache = from_cache_flags[i] if i < len(from_cache_flags) else False
                                is_serious, reason = TranslationValidator.is_serious_failure(
                                    element.full_text, translated_text, self.large_text_threshold, from_cache)
                                if is_serious or (not from_cache and reason != "翻译成功"):
                                    self._add_failed_task(element, batch_indices[i], translated_text, from_cache)
                                self._apply_translation_to_ppt_element(element, translated_text)
                        else:
                            self._add_batch_failure(batch_elements, batch_indices, translated_texts[0] if translated_texts else "批次处理失败")
                        pbar.update(len(batch_elements))
                    except Exception as e:
                        logger.error(f"批次翻译失败: {e}")
                        self._add_batch_failure(batch_elements, batch_indices, f"执行异常: {str(e)}")
                        pbar.update(len(batch_elements))

    def _print_stats(self):
        """打印统计信息（格式保持增强版）"""
        try:
            cache_stats = self.cache.stats
            print(f"\n高级PPTX翻译完成！（格式保持增强版）")
            print(f"处理时间: {self.stats['processing_time']:.1f} 秒")
            print(f"总幻灯片: {self.stats['total_slides']} 张")
            print(f"总元素: {self.stats['total_elements']} 个")
            print(f"已翻译: {self.stats['translated_elements']} 个")
            print(f"已跳过: {self.stats['skipped_elements']} 个")
            print(f"总字符数: {self.stats['total_chars']:,}")
            print(f"处理批次: {self.stats['total_batches']} 个")
            
            # PPT特有统计
            ppt_elements = []
            if self.stats['text_boxes'] > 0:
                ppt_elements.append(f"文本框: {self.stats['text_boxes']}")
            if self.stats['placeholders'] > 0:
                ppt_elements.append(f"占位符: {self.stats['placeholders']}")
            if self.stats['tables'] > 0:
                ppt_elements.append(f"表格: {self.stats['tables']}")
            if self.stats['charts'] > 0:
                ppt_elements.append(f"图表: {self.stats['charts']}")
            if self.stats['smartart'] > 0:
                ppt_elements.append(f"SmartArt: {self.stats['smartart']}")
            if self.stats['notes_pages'] > 0:
                ppt_elements.append(f"备注页: {self.stats['notes_pages']}")
            if ppt_elements:
                print(f"PPT元素: {', '.join(ppt_elements)}")
            
            print(f"API调用: {self.stats['api_calls']} 次")
            print(f"缓存命中率: {cache_stats['hit_rate']:.1%}")
            print(f"缓存节省: {self.stats['cache_savings']} 次调用")
            
            # 格式保持统计
            if self.preserve_formatting:
                total_format_attempts = self.stats['format_preserved'] + self.stats['format_fallback']
                if total_format_attempts > 0:
                    format_success_rate = self.stats['format_preserved'] / total_format_attempts
                    print(f"\n=== 格式保持统计 ===")
                    print(f"格式保持成功: {self.stats['format_preserved']} 个")
                    print(f"格式保持失败: {self.stats['format_fallback']} 个")
                    print(f"格式保持成功率: {format_success_rate:.1%}")
                    if format_success_rate >= 0.9:
                        print("✅ 格式保持效果优秀")
                    elif format_success_rate >= 0.7:
                        print("⚠️ 格式保持效果良好")
                    else:
                        print("⚠️ 格式保持效果需要改进")
            else:
                print("ℹ️ 格式保持模式已禁用")
            
            # 问题分析
            if self.stats['serious_failures'] > 0 or self.stats['minor_issues'] > 0:
                print(f"\n=== 问题分析 ===")
                if self.stats['serious_failures'] > 0:
                    print(f"严重失败: {self.stats['serious_failures']} 个")
                if self.stats['minor_issues'] > 0:
                    print(f"轻微问题: {self.stats['minor_issues']} 个")
                if self.stats['retry_attempts'] > 0:
                    print(f"\n=== 重试统计 ===")
                    print(f"重试轮数: {self.stats['retry_attempts']} 轮")
                    print(f"重试Worker: 最大{self.retry_max_workers}个，实际使用{self.stats['retry_workers_used']}个")
                    if self.stats['concurrent_retry_batches'] > 0:
                        print(f"并发重试批次: {self.stats['concurrent_retry_batches']}个")
                    print(f"缓存清理: {cache_stats['clears']} 次")
                    if self.stats['final_rescues'] > 0:
                        print(f"最终挽救: {self.stats['final_rescues']} 个")
                    if self.stats['final_failures'] > 0:
                        print(f"最终失败: {self.stats['final_failures']} 个（保持原文）")
                    else:
                        print("✅ 所有严重失败都已成功处理！")
                else:
                    print("✅ 未达到重试阈值，无需重试")
            else:
                print("✅ 无翻译问题")
            
            print(f"\n=== 配置信息 ===")
            print(f"Prompt模式: {self.stats['prompt_mode']}")
            print(f"主翻译Workers: {self.max_workers} 个")
            print(f"重试Workers: {self.retry_max_workers} 个")
            print(f"翻译备注: {'是' if self.translate_notes else '否'}")
            print(f"翻译标题: {'是' if self.translate_slide_titles else '否'}")
            print(f"格式保持: {'启用' if self.preserve_formatting else '禁用'}")
            print(f"大段文字阈值: {self.large_text_threshold} 字符")
            print(f"重试失败率阈值: {self.retry_failure_threshold:.1%}")
            
            if (self.effective_prompt_config and 
                self.effective_prompt_config.get('mode') != 'none'):
                mode = self.effective_prompt_config.get('mode', 'none')
                if mode == 'professional':
                    print(f"专业领域: {self.effective_prompt_config.get('prompt_template', 'business')}")
                elif mode == 'custom':
                    print("使用了自定义Prompt")
                # 显示增强功能
                enhancements = []
                preserve_terms = self.effective_prompt_config.get('preserve_terms')
                if preserve_terms and isinstance(preserve_terms, list):
                    enhancements.append(f"保留术语({len(preserve_terms)}个)")
                glossary = self.effective_prompt_config.get('glossary')
                if glossary and isinstance(glossary, dict):
                    enhancements.append(f"术语表({len(glossary)}条)")
                if self.effective_prompt_config.get('additional_context'):
                    enhancements.append("额外上下文")
                if enhancements:
                    print(f"增强功能: {', '.join(enhancements)}")
        except Exception as e:
            logger.warning(f"Error printing stats: {e}")

def translate_pptx_file_formatted(
    input_filepath: str,
    output_dir: str,
    target_lang: str,
    translator,
    source_lang: Optional[str] = None,
    unique_filename_base: Optional[str] = None,
    max_chunk_size: int = 6000,
    batch_size: int = 50,
    max_workers: int = 4,
    retry_max_workers: int = 3,
    reference_doc: Optional[str] = None,
    prompt_config: Optional[Dict[str, Any]] = None,
    translation_timeout: int = 45,
    max_retries: int = 6,
    large_text_threshold: int = 30,
    retry_failure_threshold: float = 0.0,
    translate_notes: bool = False,
    translate_slide_titles: bool = True,
    preserve_formatting: bool = True,  # 新增：格式保持开关
    **kwargs
) -> str:
    """
    高级格式化PPTX翻译函数（格式保持增强版）
    参数:
        translate_notes: 是否翻译备注页，默认False
        translate_slide_titles: 是否翻译幻灯片标题，默认True
        retry_max_workers: 重试时的最大worker数量，默认3
        preserve_formatting: 是否保持格式，默认True
    """
    if not os.path.exists(input_filepath):
        return f"Error: 输入文件未找到: {input_filepath}"
    os.makedirs(output_dir, exist_ok=True)
    
    # 生成输出文件路径
    input_path = Path(input_filepath)
    if unique_filename_base:
        output_filename = f"{unique_filename_base}_translated_{target_lang}.pptx"
    else:
        output_filename = f"{input_path.stem}_translated_{target_lang}.pptx"
    output_filepath = os.path.join(output_dir, output_filename)
    
    # 确保文件名唯一
    counter = 1
    while os.path.exists(output_filepath):
        if unique_filename_base:
            output_filename = f"{unique_filename_base}_translated_{target_lang}_{counter}.pptx"
        else:
            output_filename = f"{input_path.stem}_translated_{target_lang}_{counter}.pptx"
        output_filepath = os.path.join(output_dir, output_filename)
        counter += 1

    try:
        # 处理prompt配置
        effective_prompt_config = _prepare_prompt_config(prompt_config, kwargs)
        # 获取批处理设置
        batch_settings = _get_batch_settings_from_config(effective_prompt_config, kwargs)
        
        # 记录翻译开始信息
        logger.info("=== 启动高级格式化PPTX翻译（格式保持增强版）===")
        logger.info(f"输入: {os.path.basename(input_filepath)}")
        logger.info(f"目标语言: {target_lang}")
        logger.info(f"格式保持模式: {'启用' if preserve_formatting else '禁用'}")
        
        if effective_prompt_config and effective_prompt_config.get('mode') != 'none':
            mode = effective_prompt_config.get('mode', 'none')
            logger.info(f"Prompt配置: mode={mode}")
            if mode == 'professional':
                logger.info(f"专业领域: {effective_prompt_config.get('prompt_template', 'business')}")
            elif mode == 'custom':
                logger.info("使用自定义prompt")
        
        logger.info(f"PPT设置: 翻译备注={translate_notes}, 翻译标题={translate_slide_titles}")
        logger.info(f"批处理设置: batch_size={batch_settings['batch_size']}, max_chars={batch_settings['max_chars']}, workers={batch_settings['max_workers']}")
        logger.info(f"重试设置: max_retries={batch_settings['max_retries']}, timeout={batch_settings['translation_timeout']}s, retry_workers={retry_max_workers}")
        
        # 使用高级翻译器
        pptx_translator = AdvancedPptxTranslator(
            translator=translator,
            batch_size=batch_settings['batch_size'],
            max_chars=batch_settings['max_chars'],
            max_workers=batch_settings['max_workers'],
            retry_max_workers=retry_max_workers,
            prompt_config=effective_prompt_config,
            reference_doc=reference_doc,
            translation_timeout=batch_settings['translation_timeout'],
            max_retries=batch_settings['max_retries'],
            large_text_threshold=batch_settings['large_text_threshold'],
            retry_failure_threshold=batch_settings['retry_failure_threshold'],
            non_ascii_threshold=batch_settings['non_ascii_threshold'],
            translate_notes=translate_notes,
            translate_slide_titles=translate_slide_titles,
            preserve_formatting=preserve_formatting,  # 传递格式保持参数
            **kwargs
        )
        
        result_path = pptx_translator.translate_pptx(
            input_filepath=input_filepath,
            output_filepath=output_filepath,
            target_lang=target_lang,
            source_lang=source_lang
        )
        
        if not os.path.exists(result_path):
            return "Error: 输出文件未创建"
        
        file_size = os.path.getsize(result_path)
        if file_size == 0:
            return "Error: 输出文件为空"
        
        logger.info(f"高级格式化PPTX翻译成功完成！输出: {result_path}")
        return result_path
        
    except Exception as e:
        error_msg = f"高级格式化PPTX翻译失败: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return f"Error: {error_msg}"

# 测试用翻译器（格式保持增强版）
class MockPPTTranslator:
    def __init__(self):
        self.call_count = 0
        self.prompt_config = None
    def set_prompt_config(self, prompt_config):
        self.prompt_config = copy.deepcopy(prompt_config) if prompt_config else None
    def translate(self, text: str = None, target_lang: str = None, 
                 source_lang: str = None, messages: List[Dict] = None, 
                 prompt_config: Optional[Dict[str, Any]] = None,
                 config_merge_mode: str = 'merge', **kwargs) -> str:
        self.call_count += 1
        if messages:
            user_content = messages[-1]["content"]
        elif text:
            user_content = text
        else:
            return "Error: No input provided"
        # PPT特有的翻译映射（增强版）
        translation_map = {
            "标题": "Title", "内容": "Content", "项目符号": "Bullet Point",
            "图表": "Chart", "表格": "Table", "备注": "Notes",
            "幻灯片": "Slide", "演示": "Presentation", "概述": "Overview",
            "总结": "Summary", "结论": "Conclusion", "建议": "Recommendation",
            "分析": "Analysis", "数据": "Data", "趋势": "Trend",
            "方案": "Solution", "策略": "Strategy", "目标": "Goal"
        }
        # 处理presentation prompt配置
        effective_config = prompt_config or self.prompt_config
        if effective_config and effective_config.get('mode') != 'none':
            mode = effective_config.get('mode', 'none')
            if mode == 'professional':
                domain = effective_config.get('prompt_template', 'business')
                if domain == 'sales':
                    translation_map.update({
                        "销售": "Sales", "客户": "Customer", "产品": "Product",
                        "收益": "Revenue", "市场": "Market"
                    })
                elif domain == 'academic':
                    translation_map.update({
                        "研究": "Research", "论文": "Paper", "数据": "Data",
                        "方法": "Method", "结果": "Results"
                    })
        lines = user_content.split('\n')
        translated_lines = []
        for line in lines:
            if '[' in line and ']' in line:
                match = re.search(r'\[(\d+)\]\s*(.*)', line)
                if match:
                    content = match.group(2)
                    for zh, en in translation_map.items():
                        content = re.sub(r'\b' + zh + r'\b', en, content)
                    translated_lines.append(content)
        return '\n'.join(translated_lines)

if __name__ == '__main__':
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s',
        force=True
    )
    print("=== 高级格式化PPTX翻译系统（格式保持增强版）===")
    print("✅ 继承DOCX翻译器的所有核心机制")
    print("✅ 支持PPT特有元素（文本框、占位符、表格、图表、备注）")
    print("✅ 智能形状类型识别和处理")
    print("✅ 完整的格式保持系统（字体、颜色、段落格式）")
    print("✅ 智能格式分配和恢复机制")
    print("✅ 支持备注页翻译（可选）")
    print("✅ 专业演示文稿模板（business、academic、sales等）")
    print("✅ 并发重试机制，支持最大3个重试worker")
    print("✅ 智能批处理，针对PPT优化")
    print("✅ 完整的统计和监控（包含格式保持统计）")
    print("✅ 与前端完全兼容的API")
    print("✅ 可选的格式保持开关")
    
    # 演示配置
    translator = MockPPTTranslator()
    test_configs = [
        {'mode': 'professional', 'prompt_template': 'business', 'translate_notes': False, 'preserve_formatting': True},
        {'mode': 'professional', 'prompt_template': 'sales', 'translate_notes': True, 'preserve_formatting': True},
        {'mode': 'custom', 'custom_prompt': {'system': 'You are a presentation specialist...'}, 'preserve_formatting': True},
        {'mode': 'general', 'preserve_terms': ['PowerPoint', 'PPT'], 'translate_slide_titles': True, 'preserve_formatting': False},
    ]
    
    for i, config in enumerate(test_configs):
        print(f"\n=== PPT配置 {i+1}: {config.get('mode', 'default')} ===")
        mode = config.get('mode', 'none')
        print(f"模式: {mode}")
        if mode == 'professional':
            print(f"专业领域: {config.get('prompt_template', 'business')}")
        print(f"翻译备注: {config.get('translate_notes', False)}")
        print(f"翻译标题: {config.get('translate_slide_titles', True)}")
        print(f"格式保持: {config.get('preserve_formatting', True)}")
    
    print("\n=== PPT翻译系统就绪（格式保持增强版）===")
    print("支持PowerPoint演示文稿的专业翻译!")
    print("✨ 完整保持视觉格式、支持多种专业模板、智能重试机制")
    print("🎨 字体样式、颜色、段落格式完美保持")
    print("📊 详细的格式保持统计和监控")